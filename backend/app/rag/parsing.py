"""
Document Parser - extracts text from various file formats.
Supports: PDF, DOCX, TXT, MD, CSV, JSON, HTML
"""
import logging
from typing import Optional
import io
import asyncio

logger = logging.getLogger(__name__)


# File type detection
SUPPORTED_TYPES = {
    "pdf": "application/pdf",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "doc": "application/msword",
    "txt": "text/plain",
    "md": "text/markdown",
    "csv": "text/csv",
    "json": "application/json",
    "html": "text/html",
    "htm": "text/html",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}

MAX_FILE_SIZE = 10 * 1024 * 1024  # 10MB
MAX_TEXT_LENGTH = 500_000  # 500K chars max


def get_file_extension(filename: str) -> str:
    """Get file extension (lowercase, without dot)."""
    if "." not in filename:
        return ""
    return filename.rsplit(".", 1)[-1].lower()


def get_content_type(filename: str) -> str:
    """Get content type from filename."""
    ext = get_file_extension(filename)
    return SUPPORTED_TYPES.get(ext, "application/octet-stream")


def is_supported(filename: str) -> bool:
    """Check if file type is supported."""
    return get_file_extension(filename) in SUPPORTED_TYPES


async def parse_document(file_data: bytes, filename: str) -> dict:
    """
    Parse a document and extract text content.

    Returns:
        dict with: text, file_type, page_count, char_count
    """
    ext = get_file_extension(filename)

    if len(file_data) > MAX_FILE_SIZE:
        raise ValueError(f"File too large (max {MAX_FILE_SIZE // 1024 // 1024}MB)")

    try:
        if ext == "pdf":
            result = await _parse_pdf(file_data)
        elif ext in ("docx", "doc"):
            result = await _parse_docx(file_data, ext)
        elif ext == "xlsx":
            result = await _parse_xlsx(file_data)
        elif ext == "pptx":
            result = await _parse_pptx(file_data)
        elif ext in ("txt", "md", "csv", "json"):
            result = await _parse_text(file_data, ext)
        elif ext in ("html", "htm"):
            result = await _parse_html(file_data)
        else:
            # Fallback: try as text
            result = await _parse_text(file_data, "txt")

        # Truncate if too long
        if len(result["text"]) > MAX_TEXT_LENGTH:
            result["text"] = result["text"][:MAX_TEXT_LENGTH] + "\n\n[... truncated ...]"
            result["truncated"] = True

        result["file_type"] = ext
        result["char_count"] = len(result["text"])
        return result

    except Exception as e:
        logger.error(f"Failed to parse {filename}: {e}", exc_info=True)
        raise ValueError(f"Failed to parse document: {e}")


async def _parse_pdf(file_data: bytes) -> dict:
    """Parse PDF using pypdf."""
    def _sync_parse():
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(file_data))
        pages = []
        for page in reader.pages:
            text = page.extract_text() or ""
            if text.strip():
                pages.append(text.strip())
        return {
            "text": "\n\n---\n\n".join(pages),
            "page_count": len(reader.pages),
        }

    return await asyncio.to_thread(_sync_parse)


async def _parse_docx(file_data: bytes, ext: str) -> dict:
    """Parse Word document using python-docx."""
    def _sync_parse():
        from docx import Document
        doc = Document(io.BytesIO(file_data))
        paragraphs = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                paragraphs.append(text)

        # Also extract tables
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    paragraphs.append(" | ".join(cells))

        return {
            "text": "\n\n".join(paragraphs),
            "page_count": 1,  # DOCX doesn't have pages
        }

    return await asyncio.to_thread(_sync_parse)


async def _parse_xlsx(file_data: bytes) -> dict:
    """Parse Excel file using openpyxl."""
    def _sync_parse():
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(file_data), read_only=True, data_only=True)
        sheets_text = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    rows.append(" | ".join(cells))
            if rows:
                sheets_text.append(f"## Sheet: {sheet_name}\n" + "\n".join(rows))
        wb.close()
        return {
            "text": "\n\n".join(sheets_text),
            "page_count": len(wb.sheetnames),
        }

    return await asyncio.to_thread(_sync_parse)


async def _parse_pptx(file_data: bytes) -> dict:
    """Parse PowerPoint file."""
    def _sync_parse():
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_data))
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            texts.append(text)
            if texts:
                slides_text.append(f"## Slide {i}\n" + "\n".join(texts))
        return {
            "text": "\n\n".join(slides_text),
            "page_count": len(prs.slides),
        }

    return await asyncio.to_thread(_sync_parse)


async def _parse_text(file_data: bytes, ext: str) -> dict:
    """Parse plain text, markdown, CSV, JSON."""
    # Try UTF-8 first, then fallback
    try:
        text = file_data.decode("utf-8")
    except UnicodeDecodeError:
        try:
            text = file_data.decode("latin-1")
        except:
            text = file_data.decode("utf-8", errors="replace")

    return {
        "text": text,
        "page_count": 1,
    }


async def _parse_html(file_data: bytes) -> dict:
    """Parse HTML and extract text."""
    def _sync_parse():
        from bs4 import BeautifulSoup
        try:
            html = file_data.decode("utf-8")
        except UnicodeDecodeError:
            html = file_data.decode("latin-1")

        soup = BeautifulSoup(html, "html.parser")
        # Remove script and style
        for tag in soup(["script", "style"]):
            tag.decompose()
        text = soup.get_text(separator="\n", strip=True)
        return {
            "text": text,
            "page_count": 1,
        }

    return await asyncio.to_thread(_sync_parse)
